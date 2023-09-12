from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from PIL import Image


class PresentationBuilder:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Cm(29.7)
        self.prs.slide_height = Cm(21)
        self.slide_number = 0

        self.shmuts_left_text = 'Содержательный отчет к договору\n№ДЗД23-01.23 от «19» мая 2023 г.'
        self.shmuts_right_text = 'Майорова Анна Викторовна'
        self.title_mapping = {
            "за работой": "Фотографии в течение подготовки к Мероприятию",
            "на площадке": "Фотографии в течение мероприятия"
        }
        self.image_path_left = 'Снимок экрана 2023-08-22 в 13.53.33.png'
        self.image_path_right = '5.jpg'
        self.date_for_photo_left = '25.05.23\n20:19:22\nМосква'
        self.date_for_photo_right = '26.05.23\n14:29:25\nМосква'
        self.job_and_name = 'Программный менеджер 3\nМайорова Анна Викторовна'

    def create_slide(self):
        slide_layout = self.prs.slide_layouts[6]
        slide_layout.slide_width = Cm(29.7)
        slide_layout.slide_height = Cm(21)

        slide = self.prs.slides.add_slide(slide_layout)
        return slide

    def set_background_color(self, slide, color: RGBColor):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color = color
        return slide

    def black_background(self, slide):
        return self.set_background_color(slide, RGBColor(0, 0, 0))

    def grey_light_background(self, slide):
        return self.set_background_color(slide, RGBColor(217, 217, 217))

    def add_cc_logo(self, slide):
        logo_path = 'consumer_culture_logos_rgb-01-2.png'
        slide.shapes.add_picture(
            logo_path,
            Cm(-0.77),
            Cm(19.89),
            Cm(5.51),
            Cm(0.55)
        )
        return slide

    def create_time_for_photo(self, slide, date, left=False):
        left_position = Cm(12.44) if left else Cm(26.11)

        txBox = slide.shapes.add_textbox(left_position, Cm(17.55), Cm(1.84), Cm(1.42))
        text_frame = txBox.text_frame
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0

        p = text_frame.paragraphs[0]
        p.text = date

        p.font.name = 'Neue Haas Unica W1G Regular'
        p.font.size = Pt(10.5)
        p.font.color.rgb = RGBColor(255, 255, 0)
        return slide

    def add_image(self, slide, image_path, date, left=False):
        left_position = Cm(1.98) if left else Cm(15.65)

        with Image.open(image_path) as im:
            if im.size[0] > im.size[1]:
                print(f'Имя - {image_path}, слайд - номер x, фото чето надо сделать')

        slide.shapes.add_picture(image_path, left_position, Cm(2.68), Cm(12.3), Cm(16.29))
        self.create_time_for_photo(slide, date, left=left)
        return slide

    # Вспомогательная функция для создания текстового блока
    def create_textbox(self, slide, left, top, width, height, text, font_name, font_size, font_color,
                       line_spacing=None):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = txBox.text_frame
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0

        p = text_frame.paragraphs[0]
        p.text = text

        p.font.name = font_name
        p.font.size = font_size
        p.font.color.rgb = font_color
        if line_spacing:
            p.line_spacing = line_spacing

        return slide

    def create_title_photoslides(self, slide, title_photoslides_text):
        return self.create_textbox(
            slide, Cm(1.95), Cm(0.8), Cm(26.9), Cm(0.87),
            title_photoslides_text,
            'Neue Haas Unica W1G Medium', Pt(30), RGBColor(0, 0, 0)
        )

    def create_shmuts_left(self, slide):
        return self.create_textbox(
            slide, Cm(8.85), Cm(19.97), Cm(6.8), Cm(0.84),
            self.shmuts_left_text,
            'Neue Haas Unica W1G Medium', Pt(8), RGBColor(0, 0, 0), 0.9
        )

    def create_shmuts_right(self, slide):
        return self.create_textbox(
            slide, Cm(15.65), Cm(19.97), Cm(6.8), Cm(0.84),
            self.shmuts_right_text,
            'Neue Haas Unica W1G Medium', Pt(8), RGBColor(0, 0, 0), 0.9
        )

    def add_slide_numbering(self, slide):
        self.slide_number += 1

        if self.slide_number != 1:
            slide = self.create_textbox(
                slide, Cm(28.13), Cm(19.73), Cm(1), Cm(1),
                str(self.slide_number),
                'Neue Haas Unica W1G Regular', Pt(15), RGBColor(0, 0, 0)
            )
        return slide

    def title_of_shmuts(self, slide):
        # I assume job_and_name is an instance variable
        slide = self.create_textbox(
            slide, Cm(1.95), Cm(0.8), Cm(25.62), Cm(2.77),
            self.job_and_name,
            'Helvetica Neue', Pt(32), RGBColor(0, 0, 0), 1.0
        )
        return slide

    def add_photos_to_presentation(self, prs, photos, place):
        title_photoslides_text = self.title_mapping.get(place, "Заголовок по умолчанию")

        for i in range(0, len(photos), 2):
            slide = self.create_slide()
            self.create_shmuts_left(slide)
            self.create_shmuts_right(slide)
            self.create_title_photoslides(slide, title_photoslides_text)
            self.add_image(slide, photos[i], left=True)
            self.create_time_for_photo(slide, self.date_for_photo_left, left=True)

            if i + 1 < len(photos):
                self.add_image(slide, photos[i + 1], left=False)
                self.create_time_for_photo(slide, self.date_for_photo_right, left=False)

            self.add_cc_logo(slide)

        return prs

    def create_photoslide(self, image_path_left, date_left, place, image_path_right=None, date_right=None):

        title_photoslides_text = self.title_mapping.get(place, "Заголовок по умолчанию")

        slide = self.create_slide()
        self.create_shmuts_left(slide)
        self.create_shmuts_right(slide)
        self.create_title_photoslides(slide, title_photoslides_text)

        self.add_image(slide, image_path_left, date_left, left=True)

        if image_path_right:
            self.add_image(slide, image_path_right, date_right, left=False)

        self.add_cc_logo(slide)
        return slide




